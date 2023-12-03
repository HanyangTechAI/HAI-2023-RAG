import re

import docx
import requests
from kss import split_sentences
from langchain.document_loaders import PyPDFLoader
from pptx import Presentation


def normalize_text(text):
    text = re.sub(r" +", " ", text)
    text = re.sub(r"\n+", "\n", text)
    return text


def chunk_text(paragraphs, chunk_size: int = 500):
    chunks = []
    text = ""
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if len(text) + len(paragraph) >= chunk_size:
            text += "\n" + paragraph
            chunks.append(text)
            text = ""
        elif len(paragraph) > 0:
            text += paragraph + "\n"
    text = text.strip()
    if len(text) > 0:
        chunks.append(text)
    return chunks


def convert_hwp_to_txt(file_path):
    # URL for the POST request
    url = "http://localhost:40104/upload?option=all"

    # Payload for the request
    files = {"file": open(file_path, "rb")}
    data = {"option": "all"}

    # Sending the POST request
    response = requests.post(url, files=files, data=data)
    response.encoding = "utf-8"
    # Output the response
    paragraphs = split_sentences(response.text, backend="mecab")
    chunks = chunk_text(paragraphs)

    return chunks


def convert_pdf_to_txt(file_path):
    loader = PyPDFLoader(file_path)
    pages = loader.load()
    paragraphs = [page.page_content for page in pages]
    chunks = chunk_text(paragraphs)
    return chunks


def convert_docx_to_txt(file_path):
    doc = docx.Document(file_path)
    paragraphs = [paragraph.text for paragraph in doc.paragraphs]
    chunks = chunk_text(paragraphs)
    return chunks


def convert_pptx_to_txt(file_path):
    prs = Presentation(file_path)
    chunks = []

    for i, slide in enumerate(prs.slides):
        paragraphs = []
        font_sizes = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for i, paragraph in enumerate(shape.text_frame.paragraphs):
                para_text = paragraph.text.strip()
                if len(para_text) < 2:
                    continue

                font_size = 0
                for run in paragraph.runs:
                    size = run.font.size
                    size_pt = 0 if size is None else size.pt
                    if size_pt > font_size:
                        font_size = size_pt

                paragraphs.append(para_text)
                font_sizes.append(font_size)
        if len(paragraphs) == 0:
            continue

        title = None
        texts = []
        max_font_size = max(font_sizes)
        for para_text, size in zip(paragraphs, font_sizes):
            if size == max_font_size and title is None:
                title = para_text
            else:
                texts.append(para_text)
        if len(texts) == 0:
            continue

        texts = "\n".join(texts)
        texts = f"Title: {title}\n{texts}"
        chunks.append(texts)

    return chunks


def convert_file_to_txt(file_path, file_type=None):
    try:
        if file_type is None:
            file_type = file_path.strip().split(".")[-1]
        if file_type == "hwp":
            result = convert_hwp_to_txt(file_path)
        elif file_type == "pdf":
            result = convert_pdf_to_txt(file_path)
        elif file_type == "docx":
            result = convert_docx_to_txt(file_path)
        elif file_type == "pptx":
            result = convert_pptx_to_txt(file_path)
        else:
            with open(file_path, "r", encoding="utf-8") as f:
                paragraphs = f.readlines()
            result = chunk_text(paragraphs)
    except:
        result = []
    result = [normalize_text(text) for text in result]
    return result


if __name__ == "__main__":
    import glob, json

    files = glob.glob("files/*")
    print(files)

    for filepath in files:
        print(filepath)
        text = convert_file_to_txt(filepath)
        json.dump(
            {filepath: text},
            open((filepath + ".json").replace("files", "json"), "w"),
            ensure_ascii=False,
            indent=2,
        )
